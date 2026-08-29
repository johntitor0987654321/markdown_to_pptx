#!/usr/bin/env python3
"""slide.md を template.pptx のレイアウト/プレースホルダに流し込んで .pptx を生成する

見た目（色・フォント・間隔・箇条書きの記号など）は template.pptx 側のレイアウトに委ね、
本文プレースホルダには見出し/箇条書き/段落のテキストだけを流し込む。音声アイコンは
プレースホルダの下に、自前で図形として追加で配置する。画像は扱わない。

対応する記法（詳細はこのリポジトリの README.md を参照）:
  - 先頭のfrontmatter（任意）: `---` で囲んだ `key: value` 行。文字サイズ/配置はテンプレート任せ
    なので無視するが、`tts_speaker` / `tts_speed` （ナレーションのデフォルト話者・速度）だけは使う
  - スライド区切り: 独立した行の `---`
  - 見出し: スライド内で最初の `#`/`##` がタイトル、以降は小見出し
  - 箇条書き: `-`/`*`。半角スペース2つ単位のインデントでネスト
  - 太字: `**text**`
  - ナレーション: ```narration ``` コードフェンス。スライドには一切表示されない。
    `src/narration/slide-N.txt` に書き出される。フェンス内の最後の行が
    `{speaker=male speed=1.2}` の形式なら属性行として扱われ、そのスライドだけ
    話者・速度を上書きできる（speaker: male/female、speed: 1.0が標準、大きいほど速い）。
    デフォルトはfrontmatterの `tts_speaker` / `tts_speed` で指定する。
    解決結果は `src/narration/slide-N.json` に書き出され、`tts.py` が読む
  - 音声: 単独行の `[[audio: path]]`（pathはslide.mdからの相対パス）。明示指定。
    指定がなくても、ナレーションがあり `src/audio/slide-N.{wav,mp3}` が存在すれば自動で使われる。
    スピーカーアイコンとして配置され、スライド表示と同時に自動再生される（クリックでも再生できる）。
    アイコンのサイズ・位置は固定（カスタマイズ不可）

template.pptx はレイアウトを3つだけに絞ったもの（タイトルスライド / 箇条書きスライド / 左右-番号付き箇条書き）。
レイアウトの選び方:
  - デッキの1枚目、またはタイトルのみ（本文・箇条書き・音声なし）のスライド -> 「タイトルスライド - 矩形とサブタイトル」
  - それ以外 -> 「箇条書きスライド」
"""

import argparse
import json
import mimetypes
import re
import sys
from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Length

SCRIPT_DIR = Path(__file__).parent

# python-pptx は py.typed を配布しておらず内部クラスも実質非公開なので、
# prs/slide/layout/placeholder/paragraph/shape の類はすべて Any として扱う。
# block はMarkdownパース結果の辞書（"type"キーの値で形が変わる）。
Block = dict[str, Any]
Slides = list[list[Block]]
ConfigDict = dict[str, str | float]

# --- Markdownパース ---

FRONTMATTER_RE = re.compile(r"\A---\n(.*?)\n---\n?", re.DOTALL)
SLIDE_SEP_RE = re.compile(r"^\s*---\s*$", re.MULTILINE)
NARRATION_FENCE_RE = re.compile(r"^```narration[ \t]*$\n(.*?)\n^```[ \t]*$", re.MULTILINE | re.DOTALL)
ATTR_LINE_RE = re.compile(r"^\{([^{}]*)\}$")
HEADING_RE = re.compile(r"^(#{1,2})\s+(.*)$")
BULLET_RE = re.compile(r"^(\s*)[-*]\s+(.*)$")
AUDIO_RE = re.compile(r"^\[\[audio:\s*(.+?)\s*\]\]$")
BOLD_RE = re.compile(r"\*\*(.+?)\*\*")

INDENT_UNIT = 2

AUDIO_ICON = SCRIPT_DIR / "audio_icon.png"

DEFAULT_CONFIG: ConfigDict = {
    "tts_speaker": "female",
    "tts_speed": 1.0,
}


def to_float(value: str | None) -> float | None:
    """文字列をfloatに変換する。None/変換不可ならNoneを返す。"""
    if value is None:
        return None
    try:
        return float(value)
    except ValueError:
        return None


def parse_attr_tokens(text: str) -> dict[str, str]:
    """`speaker=male speed=1.2` のようなトークン列を {key: value(文字列)} にする。
    ナレーションフェンスの属性行（末尾の1行）のパースに使う。
    """
    attrs = {}
    for token in text.split():
        key, sep, value = token.partition("=")
        if sep:
            attrs[key.strip()] = value.strip()
    return attrs


def parse_frontmatter(markdown_text: str) -> tuple[ConfigDict, str]:
    """先頭のfrontmatterを読み、(設定, 残りのMarkdown) を返す。

    文字サイズ/配置系のキーはテンプレート任せなので無視し、tts_speaker/tts_speed だけ使う。
    """
    config: ConfigDict = dict(DEFAULT_CONFIG)
    match = FRONTMATTER_RE.match(markdown_text)
    if not match:
        return config, markdown_text

    for line in match.group(1).splitlines():
        key, sep, value = line.partition(":")
        key, value = key.strip(), value.strip()
        if not sep or key not in config:
            continue
        if key == "tts_speaker":
            config[key] = value
        else:
            parsed = to_float(value)
            if parsed is not None:
                config[key] = parsed

    return config, markdown_text[match.end():]


def split_slides(markdown_text: str) -> list[str]:
    """独立行の `---` でMarkdown全体をスライドごとのテキストに分割する。"""
    return [chunk.strip("\n") for chunk in SLIDE_SEP_RE.split(markdown_text)]


def resolve_narration_attrs(attrs: dict[str, str]) -> tuple[str | None, float | None]:
    """属性トークン辞書から (speaker, speed) を取り出す。"""
    return attrs.get("speaker"), to_float(attrs.get("speed"))


def extract_narration(raw: str) -> tuple[str, Block | None]:
    """narration コードフェンスを取り除いた残りのMarkdownと、そのスライドのナレーションブロック
    （無ければ None）を返す。フェンス内の最後の行が `{speaker=.. speed=..}` の形式なら属性行として
    扱い、残りの行を原稿として結合する（複数行可）。1スライドに複数フェンスがあっても最初の1つだけ使う。
    """
    match = NARRATION_FENCE_RE.search(raw)
    if not match:
        return raw, None

    lines = match.group(1).split("\n")
    speaker, speed = None, None
    if lines and ATTR_LINE_RE.match(lines[-1].strip()):
        attrs = parse_attr_tokens(ATTR_LINE_RE.match(lines.pop().strip()).group(1))
        speaker, speed = resolve_narration_attrs(attrs)

    text = "\n".join(line for line in lines if line.strip()).strip()
    remaining = raw[:match.start()] + raw[match.end():]
    block: Block = {"type": "narration", "text": text, "speaker": speaker, "speed": speed}
    return remaining, block


def parse_slide(raw: str) -> list[Block]:
    """1スライド分のMarkdown文字列を blocks のリストにする。

    block の種類:
      {"type": "title", "text": str}
      {"type": "heading", "text": str}
      {"type": "bullets", "items": [{"level": int, "text": str}]}
      {"type": "paragraph", "text": str}
      {"type": "audio", "path": str}
      {"type": "narration", "text": str, "speaker": str|None, "speed": float|None}
    """
    raw, narration_block = extract_narration(raw)

    blocks: list[Block] = []
    title_seen = False
    current_bullets: Block | None = None

    for raw_line in raw.split("\n"):
        if not raw_line.strip():
            current_bullets = None
            continue

        heading_match = HEADING_RE.match(raw_line)
        bullet_match = BULLET_RE.match(raw_line)
        audio_match = AUDIO_RE.match(raw_line.strip())

        if heading_match:
            current_bullets = None
            text = heading_match.group(2).strip()
            if not title_seen:
                blocks.append({"type": "title", "text": text})
                title_seen = True
            else:
                blocks.append({"type": "heading", "text": text})
        elif audio_match:
            current_bullets = None
            blocks.append({"type": "audio", "path": audio_match.group(1)})
        elif bullet_match:
            indent = len(bullet_match.group(1))
            level = indent // INDENT_UNIT
            text = bullet_match.group(2).strip()
            if current_bullets is None:
                current_bullets = {"type": "bullets", "items": []}
                blocks.append(current_bullets)
            current_bullets["items"].append({"level": level, "text": text})
        else:
            current_bullets = None
            blocks.append({"type": "paragraph", "text": raw_line.strip()})

    if narration_block:
        blocks.append(narration_block)

    return blocks


def parse_markdown(markdown_text: str) -> Slides:
    """Markdown全体をパースし、スライドごとのblocksのリストにする。"""
    return [parse_slide(raw) for raw in split_slides(markdown_text) if raw.strip()]


# --- ナレーションのエクスポート、音声の自動検出 ---


def export_narration(slides: Slides, config: ConfigDict) -> None:
    """ナレーションのあるスライドについて、原稿(.txt)とTTS設定(.json)を src/narration/ に書き出す。"""
    narration_dir = SCRIPT_DIR / "narration"
    for i, blocks in enumerate(slides, start=1):
        narration_block = next((b for b in blocks if b["type"] == "narration"), None)
        if narration_block is None:
            continue
        narration_dir.mkdir(exist_ok=True)
        (narration_dir / f"slide-{i}.txt").write_text(narration_block["text"] + "\n")

        tts_config = {
            "speaker": narration_block["speaker"] or config["tts_speaker"],
            "speed": narration_block["speed"] or config["tts_speed"],
        }
        (narration_dir / f"slide-{i}.json").write_text(
            json.dumps(tts_config, ensure_ascii=False, indent=2) + "\n"
        )


def find_auto_audio(slide_index: int) -> str | None:
    """src/audio/slide-N.wav（または.mp3）があれば、slide.mdからの相対パスを返す。"""
    for ext in ("wav", "mp3"):
        if (SCRIPT_DIR / "audio" / f"slide-{slide_index}.{ext}").exists():
            return f"src/audio/slide-{slide_index}.{ext}"
    return None


# --- テンプレートへの流し込み ---

LAYOUT_TITLE = "タイトルスライド - 矩形とサブタイトル"
LAYOUT_CONTENT = "箇条書きスライド"
SLIDE_NUMBER_IDX = 12

BODY_MARGIN_X = Inches(0.7)
DEFAULT_AUDIO_ICON_SIZE = Inches(0.5)


def remove_all_slides(prs: Any) -> None:
    """テンプレートに元から入っている実データのスライドを全部消し、レイアウト（マスター）だけ残す。"""
    xml_slides = prs.slides._sldIdLst
    for sld in list(xml_slides):
        r_id = sld.get(qn("r:id"))
        prs.part.drop_rel(r_id)
        xml_slides.remove(sld)


def remove_decorative_lines(prs: Any) -> None:
    """レイアウトに入っている非プレースホルダのLINEシェイプ（タイトルと本文の区切り線など）を取り除く。"""
    for layout in prs.slide_layouts:
        for shape in list(layout.shapes):
            if not shape.is_placeholder and shape.shape_type == MSO_SHAPE_TYPE.LINE:
                shape._element.getparent().remove(shape._element)


def add_slide_number(slide: Any, layout: Any) -> None:
    """ページ番号プレースホルダをレイアウトから複製する（add_slide()では自動複製されないため）。"""
    if any(ph.placeholder_format.idx == SLIDE_NUMBER_IDX for ph in slide.placeholders):
        return

    layout_ph = next(
        (ph for ph in layout.placeholders if ph.placeholder_format.idx == SLIDE_NUMBER_IDX), None
    )
    if layout_ph is None:
        return

    sp = deepcopy(layout_ph._element)
    new_id = max((shape.shape_id for shape in slide.shapes), default=0) + 1
    sp.find(qn("p:nvSpPr")).find(qn("p:cNvPr")).set("id", str(new_id))
    slide.shapes._spTree.append(sp)


def layout_by_name(prs: Any, name: str) -> Any:
    """名前でスライドレイアウトを探す。無ければ例外を投げる。"""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"レイアウトが見つかりません: {name}")


def add_bold_runs(paragraph: Any, text: str) -> None:
    """フォントサイズはテンプレート側のデフォルトに任せ、太字だけ上書きする。"""
    pos = 0
    for match in BOLD_RE.finditer(text):
        if match.start() > pos:
            paragraph.add_run().text = text[pos:match.start()]
        run = paragraph.add_run()
        run.text = match.group(1)
        run.font.bold = True
        pos = match.end()
    if pos < len(text):
        paragraph.add_run().text = text[pos:]


def choose_layout(prs: Any, blocks: list[Block], is_first_slide: bool) -> Any:
    """1枚目、または本文が無いスライドはタイトルレイアウト、それ以外は箇条書きレイアウトを選ぶ。"""
    has_body = any(b["type"] in ("heading", "paragraph", "bullets", "audio") for b in blocks)
    if is_first_slide or not has_body:
        return layout_by_name(prs, LAYOUT_TITLE)
    return layout_by_name(prs, LAYOUT_CONTENT)


def fill_body_placeholder(placeholder: Any, blocks: list[Block]) -> None:
    """見出し/段落/箇条書きのテキストを本文プレースホルダに流し込む。"""
    tf = placeholder.text_frame
    first = True
    for block in blocks:
        if block["type"] == "heading":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            add_bold_runs(p, block["text"])
            for run in p.runs:
                run.font.bold = True
        elif block["type"] == "paragraph":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            add_bold_runs(p, block["text"])
        elif block["type"] == "bullets":
            for item in block["items"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.level = min(item["level"], 8)
                add_bold_runs(p, item["text"])


def add_audio_icon(slide: Any, audio_path: Path, left: Length, top: Length, size: Length) -> Any:
    """クリック/自動再生される音声アイコンを配置する。

    python-pptx には音声専用のAPIがない（動画用の add_movie() のみ）。add_movie() は内部で
    メディアパーツの追加・関係付け・クリック再生用の hlinkClick まで正しく作ってくれるので、
    それを流用する。要素タグ <a:videoFile> を <a:audioFile> に、参照している関係(rel)の
    Type も video から audio に付け替える（両方揃えないとGoogleスライド等で再生できない）。
    """
    mime_type = mimetypes.guess_type(str(audio_path))[0] or "audio/mpeg"
    shape = slide.shapes.add_movie(
        str(audio_path), left, top, size, size,
        poster_frame_image=str(AUDIO_ICON) if AUDIO_ICON.exists() else None,
        mime_type=mime_type,
    )
    video_file = shape._element.find(".//" + qn("a:videoFile"))
    video_file.tag = qn("a:audioFile")

    r_id = video_file.get(qn("r:link"))
    rel = slide.part.rels[r_id]
    # rel.reltype は読み取り専用の lazyproperty で add_movie() 実行中に "video" としてキャッシュ
    # 済みのため、通常の属性代入では書き換えられない。__dict__ を直接上書きして迂回する。
    rel.__dict__["reltype"] = RELATIONSHIP_TYPE.AUDIO

    return shape


def set_autoplay(slide: Any) -> None:
    """そのスライドの音声を、クリック待ちではなくスライド表示と同時に再生させる。

    add_movie() が組み立てるクリック再生用の <p:timing> のトリガー条件
    （delay="indefinite" = イベント待ち）を "0"（即時）に変え、<p:video> タグも
    <p:audio> にする（<a:videoFile>のときと同種のタグ/実体の不整合を避けるため）。
    """
    timing = slide._element.find(qn("p:timing"))
    if timing is None:
        return

    for video_node in timing.findall(".//" + qn("p:video")):
        video_node.tag = qn("p:audio")

    for cond in timing.findall(".//" + qn("p:cond")):
        if cond.get("delay") == "indefinite":
            cond.set("delay", "0")


def fill_title_layout(title_placeholder: Any, title_text: str, blocks: list[Block]) -> None:
    """タイトルレイアウト用: このテンプレートはタイトル用プレースホルダ1つに
    複数行（タイトル+サブタイトル相当の段落）を入れる作りだった。"""
    tf = title_placeholder.text_frame
    tf.text = title_text
    for block in blocks:
        if block["type"] == "paragraph":
            p = tf.add_paragraph()
            add_bold_runs(p, block["text"])


def resize_body_for_audio(body_ph: Any) -> Length:
    """本文と音声アイコンが同居するスライド用に、本文プレースホルダを上側45%へ縮める。
    残り(下側55%)を音声アイコン用に空け、そこの開始位置(top)を返す。

    プレースホルダは元々レイアウト側の位置/サイズを継承しており明示的な xfrm を
    持たないことが多いため、heightだけ設定すると width/left/top が0のxfrmが新規に
    作られてしまい、見えないプレースホルダになる。4つとも明示的に設定する。
    """
    left, top, width = body_ph.left, body_ph.top, body_ph.width
    bottom = top + body_ph.height
    new_height = int(bottom - top) * 45 // 100

    body_ph.left = left
    body_ph.top = top
    body_ph.width = width
    body_ph.height = new_height
    return top + new_height + Inches(0.15)


def place_body_text(
    prs: Any, body_ph: Any | None, text_blocks: list[Block], audio_blocks: list[Block]
) -> tuple[Length, Length, Length]:
    """本文プレースホルダにテキストを流し込み、音声アイコンの配置開始位置
    (content_left, content_width, top) を返す。"""
    if body_ph is None:
        # このレイアウトに本文プレースホルダが無い場合のフォールバック（通常は発生しない想定）
        return BODY_MARGIN_X, prs.slide_width - 2 * BODY_MARGIN_X, Inches(1.5)

    content_left = body_ph.left
    content_width = body_ph.width
    content_bottom = body_ph.top + body_ph.height

    if text_blocks and audio_blocks:
        top = resize_body_for_audio(body_ph)
        fill_body_placeholder(body_ph, text_blocks)
    elif text_blocks:
        fill_body_placeholder(body_ph, text_blocks)
        top = content_bottom
    else:
        top = body_ph.top

    return content_left, content_width, top


def place_audio_icons(
    slide: Any, audio_blocks: list[Block], base_dir: Path, content_left: Length, content_width: Length, top: Length
) -> None:
    """音声ブロックを上から順に配置する（サイズ・位置は固定）。ファイルが無ければ代わりにエラーメッセージを置く。"""
    for block in audio_blocks:
        audio_path = (base_dir / block["path"]).resolve()
        if not audio_path.exists():
            box = slide.shapes.add_textbox(content_left, top, content_width, Inches(0.4))
            box.text_frame.text = f"[音声が見つかりません: {block['path']}]"
            top += Inches(0.5)
            continue

        add_audio_icon(slide, audio_path, content_left, top, DEFAULT_AUDIO_ICON_SIZE)
        top += DEFAULT_AUDIO_ICON_SIZE + Inches(0.2)


def build_slide(prs: Any, blocks: list[Block], base_dir: Path, is_first_slide: bool) -> None:
    """1スライド分のblocksをテンプレートのレイアウトに流し込む。"""
    # レイアウトを選び、スライドを追加する
    layout = choose_layout(prs, blocks, is_first_slide)
    slide = prs.slides.add_slide(layout)
    add_slide_number(slide, layout)
    placeholders = {ph.placeholder_format.idx: ph for ph in slide.placeholders}

    title_block = next((b for b in blocks if b["type"] == "title"), None)
    title_text = title_block["text"] if title_block else ""

    # タイトルレイアウトは専用の1プレースホルダに全部詰めるだけで完結する
    if layout.name == LAYOUT_TITLE:
        fill_title_layout(placeholders[0], title_text, blocks)
        return

    if 0 in placeholders:
        placeholders[0].text_frame.text = title_text

    # 本文（見出し/段落/箇条書き）と音声を分けて配置する
    text_blocks = [b for b in blocks if b["type"] in ("heading", "paragraph", "bullets")]
    audio_blocks = [b for b in blocks if b["type"] == "audio"]

    content_left, content_width, top = place_body_text(prs, placeholders.get(1), text_blocks, audio_blocks)
    place_audio_icons(slide, audio_blocks, base_dir, content_left, content_width, top)

    set_autoplay(slide)


def attach_auto_audio(blocks: list[Block], slide_index: int) -> None:
    """`[[audio: path]]` の明示指定が無く、ナレーションだけあるスライドに、
    自動検出した src/audio/slide-N.* を音声ブロックとして追加する。"""
    has_audio = any(b["type"] == "audio" for b in blocks)
    has_narration = any(b["type"] == "narration" for b in blocks)
    if has_audio or not has_narration:
        return

    auto_path = find_auto_audio(slide_index)
    if auto_path:
        blocks.append({"type": "audio", "path": auto_path})


def build_pptx(template_path: Path, slides: Slides) -> Any:
    """全スライド分のblocksからPresentationを組み立てる。"""
    # テンプレートを読み込み、実データのスライドを除いてレイアウトだけにする
    prs = Presentation(str(template_path))
    remove_all_slides(prs)
    remove_decorative_lines(prs)

    # スライドごとに音声を自動検出しつつ、テンプレートへ流し込む
    for i, blocks in enumerate(slides, start=1):
        attach_auto_audio(blocks, i)
        build_slide(prs, blocks, base_dir=template_path.parent, is_first_slide=(i == 1))

    return prs


def parse_args() -> argparse.Namespace:
    """コマンドライン引数（slide.mdのパス、--export-narration）を定義してパースする。"""
    parser = argparse.ArgumentParser(description="slide.md を .pptx に変換する")
    parser.add_argument("slide_md", help="slide.mdのパス")
    parser.add_argument(
        "--export-narration",
        action="store_true",
        help="ナレーション原稿（src/narration/*.txt, *.json）だけ書き出し、pptxは生成しない",
    )
    return parser.parse_args()


def main() -> None:
    """slide.mdを読み、ナレーションを書き出し、（--export-narrationで無ければ）pptxを生成する。"""
    args = parse_args()

    # Markdownをパースし、ナレーションを書き出す（tts.pyの入力になる）
    src = Path(args.slide_md).resolve()
    config, markdown_text = parse_frontmatter(src.read_text())
    slides = parse_markdown(markdown_text)
    export_narration(slides, config=config)

    if args.export_narration:
        # 音声合成(tts.py)がまだなので、ここではpptxを作らない
        print(f"wrote narration for {len(slides)} slides")
        return

    # 事前条件チェック: template.pptx が無ければここで終了する
    template_path = src.parent / "template.pptx"
    if not template_path.exists():
        print(f"template.pptx が見つかりません: {template_path}")
        sys.exit(1)

    # テンプレートに流し込んで.pptxを組み立て、保存する
    prs = build_pptx(template_path, slides)
    dst = src.with_suffix(".pptx")
    prs.save(str(dst))
    print(f"wrote {dst} ({len(slides)} slides)")


if __name__ == "__main__":
    main()
